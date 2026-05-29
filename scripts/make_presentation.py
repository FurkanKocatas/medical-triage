"""Türkçe Tıbbi Triaj — Final sunumu (PPTX) üretici.

Kullanım:
    py -3.14 scripts/make_presentation.py
Çıktı:
    sunum.pptx  (proje kökünde)
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
PLOTS = ROOT / "eval" / "plots"

# --- Palet ---
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
GREY = RGBColor(0x55, 0x5F, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x2A, 0x33)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
YELLOW = RGBColor(0xF9, 0xA8, 0x25)
RED = RGBColor(0xC6, 0x28, 0x28)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False, font="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def header(slide, title, kicker=None):
    """Üst başlık bandı."""
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Inches(0.06), ACCENT)
    tf = textbox(slide, Inches(0.5), Inches(0.18), SW - Inches(1.0), Inches(0.9),
                 anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker + "\n"; _set(r, 12, RGBColor(0x9F, 0xB3, 0xC8), bold=True)
    r = p.add_run(); r.text = title; _set(r, 26, WHITE, bold=True)


def footer(slide, idx):
    tf = textbox(slide, Inches(0.4), SH - Inches(0.45), Inches(9), Inches(0.35))
    r = tf.paragraphs[0].add_run()
    r.text = "Türkçe Tıbbi Triaj • SLM Benchmark • Selçuk Ü. 2026"
    _set(r, 9, GREY)
    tf2 = textbox(slide, SW - Inches(1.2), SH - Inches(0.45), Inches(0.8), Inches(0.35))
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r2 = tf2.paragraphs[0].add_run(); r2.text = str(idx); _set(r2, 11, NAVY, bold=True)


def bullets(slide, items, x=Inches(0.6), y=Inches(1.5), w=None, h=None, size=18, gap=8):
    w = w or (SW - Inches(1.2))
    h = h or (SH - Inches(2.2))
    tf = textbox(slide, x, y, w, h)
    first = True
    for it in items:
        lvl = 0
        color = DARK
        bold = False
        if isinstance(it, tuple):
            text, lvl = it[0], it[1]
            if len(it) > 2:
                color = it[2]
            if len(it) > 3:
                bold = it[3]
        else:
            text = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        bullet = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = bullet + text
        _set(r, size - lvl * 2, color, bold=bold or (lvl == 0))
    return tf


def content(title, items, kicker=None, idx=0, size=18):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, LIGHT)
    header(s, title, kicker)
    bullets(s, items, size=size)
    footer(s, idx)
    return s


def image_slide(title, img, items, kicker=None, idx=0, size=15):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, LIGHT)
    header(s, title, kicker)
    # sol metin, sağ görsel
    bullets(s, items, x=Inches(0.5), y=Inches(1.5), w=Inches(5.6),
            h=SH - Inches(2.2), size=size, gap=6)
    p = PLOTS / img
    if p.exists():
        # sağ yarıya sığdır
        s.shapes.add_picture(str(p), Inches(6.4), Inches(1.5),
                             height=Inches(5.0))
    footer(s, idx)
    return s


def table_slide(title, headers, rows, kicker=None, idx=0, note=None,
                col_widths=None, font=11, highlight_row=None):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, LIGHT)
    header(s, title, kicker)
    nrows, ncols = len(rows) + 1, len(headers)
    left, top = Inches(0.45), Inches(1.55)
    width = SW - Inches(0.9)
    height = Inches(0.5) * nrows
    gtable = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = Inches(cw)
    # header satırı
    for j, htxt in enumerate(headers):
        c = gtable.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = c.text_frame.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
        r = para.add_run(); r.text = htxt; _set(r, font, WHITE, bold=True)
    # gövde
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gtable.cell(i, j)
            c.fill.solid()
            if highlight_row is not None and i == highlight_row:
                c.fill.fore_color.rgb = RGBColor(0xFD, 0xEC, 0xEA)
            else:
                c.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xEC, 0xF1, 0xF6)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = c.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = para.add_run(); r.text = str(val)
            _set(r, font, DARK, bold=(highlight_row is not None and i == highlight_row))
    if note:
        tf = textbox(s, Inches(0.5), top + height + Inches(0.15), SW - Inches(1.0), Inches(1.2))
        r = tf.paragraphs[0].add_run(); r.text = note; _set(r, 13, GREY, italic=True)
    footer(s, idx)
    return s


# ============================== SLAYTLAR ==============================

# 1) Başlık
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.08), ACCENT)
# triaj renk şeritleri
rect(s, 0, 0, Inches(0.25), SH, GREEN)
rect(s, Inches(0.25), 0, Inches(0.25), SH, YELLOW)
rect(s, Inches(0.5), 0, Inches(0.25), SH, RED)
tf = textbox(s, Inches(1.2), Inches(1.7), SW - Inches(2.2), Inches(3.0))
r = tf.paragraphs[0].add_run()
r.text = "Türkçe Tıbbi Triaj ve Bölüm Yönlendirme için\nKüçük Dil Modellerinin İnce Ayarı"
_set(r, 36, WHITE, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(14)
r = p.add_run(); r.text = "SmolLM2-360M · Qwen2.5-1.5B · Gemma-4-E2B — Karşılaştırmalı Benchmark"
_set(r, 18, RGBColor(0x9F, 0xB3, 0xC8))
tf2 = textbox(s, Inches(1.2), Inches(4.9), SW - Inches(2.2), Inches(2.0))
for line, sz in [("Ebubekir Bayar (258273002001)  ·  Furkan Kocataş (258273002008)", 16),
                 ("Selçuk Üniversitesi — Fen Bilimleri Enstitüsü, Bilgisayar Mühendisliği (Y.L.)", 13),
                 ("Büyük Dil Modelleri Dersi — Final Projesi · Konya, 2026", 13)]:
    p = tf2.add_paragraph(); r = p.add_run(); r.text = line
    _set(r, sz, WHITE if sz > 14 else RGBColor(0xC8, 0xD2, 0xDC))

# 2) Problem & Motivasyon
content("Problem Tanımı ve Motivasyon", [
    ("Türkiye'de 2023'te acil servise yıllık 130M+ başvuru — nüfusun ~1,5 katı, OECD ortalamasının çok üzerinde.", 0),
    ("Yoğunluk → triaj hemşiresi başına vaka artışı, karar süresinde sıkışma; hatalı bölüm yönlendirmesi gereksiz konsültasyon/tetkik.", 0),
    ("Hedef: Türkçe serbest semptom metninden eş-zamanlı iki çıktı.", 0),
    ("Triaj aciliyet: Yeşil / Sarı / Kırmızı", 1, GREY),
    ("Bölüm yönlendirme: Kardiyoloji, Nöroloji, Dahiliye, …", 1, GREY),
    ("Neden Küçük Dil Modeli (SLM)? GPT-4/Med-PaLM gibi kapalı modeller 3 engelle karşılaşır:", 0),
    ("KVKK / veri yerelliği (veri yurt dışına çıkamaz)", 1, GREY),
    ("Çağrı başına API maliyeti — ölçekte bütçe baskısı", 1, GREY),
    ("Çevrimdışı / gerçek zamanlı çalışma ihtiyacı", 1, GREY),
    ("Klinik güvenlik metrikleri ön planda: False Green Rate (FGR), Kırmızı için FNR-Red.", 0, ACCENT, True),
], kicker="1 / Giriş", idx=2, size=17)

# 3) İlgili Çalışmalar
content("İlgili Çalışmalar (Literatür Özeti)", [
    ("Klasik triaj sistemleri: Manchester Triage System (MTS), Emergency Severity Index (ESI).", 0),
    ("NLP tabanlı triaj: TF-IDF → BioBERT / ClinicalBERT / PubMedBERT; serbest metinden 0,84 AUROC'a ulaşan çalışmalar.", 0),
    ("LLM dönemi: Med-PaLM ve Med-Gemini USMLE'de uzman seviyesi; açık kaynak Llama-2/Mistral üzerine LoRA ince ayar (2024+).", 0),
    ("Türkçe medikal NLP: İngilizce muadillere göre 5–8 puan F1 düşüşü; kaynak ve etiketli veri kısıtlı.", 0),
    ("Boşluk: 3 SLM ailesinin Türkçe triajda, klinik güvenlik metrikleriyle yan yana değerlendirildiği kamuya açık bir benchmark yok.", 0, ACCENT, True),
    ("Bu çalışma bu boşluğu hedefler (15+ kaynaklı sistematik tarama — bkz. rapor).", 0),
], kicker="2 / Literatür", idx=3, size=17)

# 4) Veri & Etiketleme
image_slide("Yöntem — Veri ve Etiketleme", "class_distribution.png", [
    ("Kaynak: alibayram/doktorsitesi (HF, gated) → 5.000 örnek; seed=42.", 0),
    ("116 alt-uzmanlık → 23 ana bölüme normalize.", 0),
    ("Etiketleme: Kaggle triaj veri setinden çıkarılan ayırt edici anahtar kelimelerle 3 katmanlı KURAL TABANLI sınıflandırıcı.", 0),
    ("STRONG_RED + akut başlangıç + geçmiş-zaman filtresi.", 1, GREY),
    ("Ciddi dengesizlik: Kırmızı yalnızca %4,1.", 0, ACCENT, True),
    ("80/10/10 stratified (triaj×bölüm): 3.972 / 475 / 553.", 0),
], kicker="3 / Yöntem", idx=4, size=15)

# 5) Modeller
table_slide("Yöntem — Model Seçimi (Zorunlu Setten 3)",
            ["Model", "Parametre", "Bağlam", "Tokenizer", "Mimari"],
            [["SmolLM2-360M", "362 M", "8K", "49.152", "Llama2-vari, RoPE"],
             ["Qwen2.5-1.5B", "1,5 B", "32K", "151.643", "GQA + SwiGLU, çok dilli"],
             ["Gemma-4-E2B", "~2 B*", "8K", "256.000", "MoE-vari, edge"]],
            kicker="3 / Yöntem", idx=5, font=14,
            col_widths=[2.6, 1.8, 1.4, 2.0, 4.6],
            note="* Gemma-4-E2B toplam ağırlık ~5,1 B; etkin parametre ~2 B. "
                 "Üç farklı mimari ailesi (kapasite tavanı, çok dillilik, edge MoE) bilinçli seçildi.")

# 6) İnce Ayar Stratejisi
content("Yöntem — İnce Ayar Stratejisi", [
    ("QLoRA (4-bit NF4 + bf16 compute): taban ağırlıklar 4-bit, düşük VRAM.", 0),
    ("LoRA bf16 (no 4-bit): taban ağırlıklar bf16; 12 GB VRAM'de Qwen için mümkün.", 0),
    ("LoRA bf16 + Dengeli veri: Kırmızı sınıfı 4× temiz oversample (train_balanced).", 0),
    ("Ortak LoRA: r=16, alpha=32, dropout=0,05 (Gemma r=8); paged_adamw_8bit, gradient checkpointing.", 0, GREY),
    ("Prompt şablonu (3 alanlı):  ### Görev → ### Semptom → ### Cevap (JSON: triaj, bölüm, neden).", 0),
    ("Donanım: ilk faz NVIDIA DGX Spark GB10; erişim kesintisi sonrası RTX 4070 12 GB (Windows). Latency RTX 4070'te.", 0, GREY),
], kicker="4 / İnce Ayar", idx=6, size=17)

# 7) Eğitim Süreç Metrikleri
table_slide("Eğitim Süreç Detayları (Her Model)",
            ["Model / Yöntem", "Epoch", "LR", "bs×ga", "max_len", "Süre", "eval_loss"],
            [["SmolLM2 — QLoRA r16", "3", "3e-4", "8×2", "768", "62,3 dk", "1,625"],
             ["Qwen2.5 — QLoRA r16", "2", "2e-4", "4×4", "896", "89,2 dk", "1,662"],
             ["Qwen2.5 — LoRA bf16", "2", "2e-4", "2×8", "896", "35,6 dk", "1,615"],
             ["Qwen2.5 — LoRA dengeli", "2", "2e-4", "2×8", "896", "~36 dk", "—"],
             ["Gemma-4 — QLoRA r8", "1", "1e-4", "1×16", "512", "42,8 dk", "—"]],
            kicker="5 / Sonuçlar", idx=7, font=13,
            col_widths=[3.4, 1.1, 1.2, 1.3, 1.4, 1.6, 1.5],
            note="LoRA bf16, QLoRA'ya göre ~2,5× daha hızlı ve daha düşük eval_loss "
                 "(12 GB'de 1,5 B model için QLoRA gerekli değil).")

# 8) Zero-shot Baseline
content("Zero-Shot Baseline (Fine-tuning Öncesi)", [
    ("Üç modelin de zero-shot performansı aynı metriklerle ölçüldü (n=553).", 0),
    ("SmolLM2: JSON parse başarısı yalnızca %2,7 → accuracy ≈ 0.", 0, GREY),
    ("Qwen2.5: parse %80 ama geçerli Türkçe etiket üretemiyor → accuracy = 0.", 0, GREY),
    ("Gemma-4: parse %56, accuracy 0,139 (Macro-F1 0,118).", 0, GREY),
    ("Sonuç: Görev-özel çıktı formatı (JSON + Türkçe etiket) zero-shot ile güvenilir üretilemiyor.", 0, ACCENT, True),
    ("→ İnce ayar bu görev için zorunlu; baseline ile FT arasında büyük sıçrama var.", 0),
], kicker="5 / Sonuçlar", idx=8, size=17)

# 9) Standart Sonuç Tablosu (ANA)
table_slide("Standart Sonuç Tablosu (Benchmark Protokolü)",
            ["Model", "Param", "FT", "Acc", "Macro-F1", "Eğitim", "Inf. (ms)", "Boyut (MB)"],
            [["SmolLM2-360M", "362M", "QLoRA r16", "0,562", "0,244", "62,3 dk", "1.706", "21,0"],
             ["Qwen2.5-1.5B", "1,5B", "QLoRA r16", "0,622", "0,260", "89,2 dk", "1.391", "48,4"],
             ["Qwen2.5-1.5B", "1,5B", "LoRA bf16", "0,624", "0,256", "35,6 dk", "1.463", "85,3"],
             ["Gemma-4-E2B", "~2B", "QLoRA r8", "0,624", "0,256", "42,8 dk", "2.733", "70,3"],
             ["Qwen2.5-1.5B", "1,5B", "LoRA dengeli", "0,618", "0,390", "~36 dk", "1.407", "85,3"]],
            kicker="5 / Sonuçlar", idx=9, font=12.5, highlight_row=5,
            col_widths=[2.4, 1.0, 1.7, 1.0, 1.3, 1.4, 1.3, 1.5],
            note="Boyutlar adapter-only. En iyi: dengeli Qwen2.5-1.5B (Macro-F1 0,390). "
                 "Metrikler: Accuracy + Macro-F1 (Weighted-F1 ve güvenlik metrikleri raporda).")

# 10) Sınıf Dengesizliği Müdahalesi
image_slide("Sınıf Dengesizliği Müdahalesi (Qwen2.5-1.5B)", "metric_comparison.png", [
    ("İlk FT'de tüm modeller ~%62,4'e yakınsadı = test Sarı oranı → hepsi Sarı'ya çöktü, FNR-Red = 1,000.", 0),
    ("Kırmızı 4× oversample (dengeli veri) sonrası:", 0),
    ("Macro-F1: 0,256 → 0,390  (+%52 göreli)", 1, GREEN, True),
    ("FNR-Red: 1,000 → 0,529  (−0,471)", 1, RED, True),
    ("FGR: 0,000 korundu (yanlış-yeşil yok).", 1, GREY),
    ("Yine de FNR-Red, klinik eşiğin (<%5–10) çok üzerinde.", 0, ACCENT, True),
], kicker="5 / Sonuçlar", idx=10, size=15)

# 11) Pareto + Cascade
image_slide("Boyut–Performans (Pareto) ve Cascade Bulgusu", "pareto_scatter.png", [
    ("Pareto: Qwen2.5-1.5B en iyi boyut/performans dengesi.", 0),
    ("Gemma ~2B, Qwen 1,5B ile aynı Macro-F1 (0,256) — kapasite tek başına yetmiyor.", 0, GREY),
    ("Cascade (L1→L2→L3) NEGATİF bulgu:", 0, ACCENT, True),
    ("Tek model Qwen-dengeli ile ~bit-aynı sonuç.", 1, GREY),
    ("~11× ek gecikme, katma değer yok.", 1, GREY),
    ("Neden: erken durdurma çalışmıyor — hiçbir katman Yeşil üretemiyor, her örnek L3'e eskale oluyor.", 1, GREY),
], kicker="5 / Sonuçlar", idx=11, size=15)

# 12) Hata Analizi
image_slide("Hata Analizi (Qwen2.5 Dengeli)", "confusion_matrices_all.png", [
    ("211/553 hata. Baskın patern: Yeşil→Sarı (164, %77,7).", 0),
    ("Vaka 1 (idx=3): geçmiş travma → PAST_TENSE filtresi Kırmızı'yı maskeledi (Kırmızı→Sarı).", 0, GREY),
    ("Vaka 2 (idx=12): kronik/danışma sorusu → Yeşil→Sarı yumuşak yüksek tahmin.", 0, GREY),
    ("Vaka 3 (idx=302): cascade aşırı eskalasyon → Yeşil→Kırmızı.", 0, GREY),
    ("Vaka 4 (idx=10): meme asimetrisi → bölüm Genel Cerrahi yerine Ortopedi (anatomik karışıklık).", 0, GREY),
    ("Vaka 5 (idx=15): kritik vaka doğru Kırmızı yakalandı (bölüm Kardiyoloji ~ kabul edilebilir).", 0, GREEN),
    ("Yeşil sınıfı 9 yapılandırmanın hiçbirinde yakalanamadı (F1=0): dilsel örtüşme + etiket gürültüsü + çoğunluk baskısı.", 0, ACCENT, True),
], kicker="6 / Hata Analizi", idx=12, size=13)

# 13) Tartışma
content("Tartışma", [
    ("Boyut-performans: Qwen (1,5B) > SmolLM2 (360M) +0,146 Macro-F1; ama Gemma (~2B) = Qwen — veri ve veriye uyarlanmış mimari, kapasite kadar önemli.", 0),
    ("LoRA bf16 vs QLoRA: bf16 2,5× hızlı + daha düşük eval_loss → 12 GB'de QLoRA gereksiz; QLoRA'nın değeri <8 GB veya 7B+ modellerde.", 0),
    ("Dengeleme: veri-seviyesi oversample tek başına yetmez; loss-seviyesi (class-weighted CE, focal loss) gerekli.", 0),
    ("Edge uygunluk:", 0),
    ("SmolLM2 (0,87 GB) Raspberry Pi sınıfı — ama Macro-F1 yetersiz.", 1, GREY),
    ("Qwen2.5 (3,28 GB) telefon/Jetson Orin Nano — kabul edilebilir başlangıç.", 1, GREEN),
    ("Gemma (10,42 GB) yalnızca masaüstü GPU/bulut — edge vaadi Türkçe'de gerçekleşmedi.", 1, GREY),
], kicker="7 / Tartışma", idx=13, size=16)

# 14) Sınırlamalar
content("Sınırlamalar", [
    ("FNR-Red = 0,529 — klinik eşiğin (<%5–10) çok üzerinde; üretime uygun DEĞİL.", 0, ACCENT, True),
    ("Yeşil sınıfı F1 = 0 (hiç yakalanamadı).", 0),
    ("Tek tohum (seed=42): çoklu tohum (≥3) ile ortalama±std ve anlamlılık testi yapılmadı → gelecek çalışma.", 0),
    ("Kural tabanlı etiketler doktor doğrulamasından geçmedi (etiket gürültüsü).", 0),
    ("doktorsitesi dağılımı kadın doğum/pediatri ağırlıklı; Acil Tip az temsil.", 0, GREY),
    ("Sistem bir ARAŞTIRMA PROTOTİPİDİR — hekim doğrulaması olmadan kullanılmamalıdır.", 0, ACCENT, True),
], kicker="8 / Sınırlamalar", idx=14, size=17)

# 15) Sonuç & Gelecek
content("Sonuç ve Gelecek Çalışmalar", [
    ("3 SLM ailesi Türkçe triajda klinik güvenlik metrikleriyle ilk kez yan yana kıyaslandı.", 0),
    ("En iyi: dengeli Qwen2.5-1.5B (Macro-F1 0,390; FNR-Red 0,529) — açık kaynakla ulaşılabilir başlangıç sınırı.", 0),
    ("Gelecek çalışmalar:", 0),
    ("Class-weighted CE + focal loss (Kırmızı recall hedefi %70–80).", 1, GREY),
    ("İki aşamalı Yeşil-vs-rest + back-translation ile sentetik Kırmızı.", 1, GREY),
    ("≥3 tohumla varyans + istatistiksel anlamlılık; TF-IDF/mBERT baseline.", 1, GREY),
    ("Doktor doğrulamalı ≥2.000 test seti; HF Hub yayını; ortak benchmark makalesi.", 1, GREY),
], kicker="9 / Sonuç", idx=15, size=16)

# 16) Demo + Teşekkür
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.7), SW, Inches(0.06), ACCENT)
tf = textbox(s, Inches(1.2), Inches(1.6), SW - Inches(2.4), Inches(2.2), anchor=MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run(); r.text = "Canlı Demo & Sorular"; _set(r, 40, WHITE, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(16)
r = p.add_run()
r.text = ("Demo: tek satır Türkçe semptom → {triaj, bölüm, neden} JSON\n"
          "(dengeli Qwen2.5-1.5B + cascade)")
_set(r, 18, RGBColor(0x9F, 0xB3, 0xC8))
tf2 = textbox(s, Inches(1.2), Inches(4.1), SW - Inches(2.4), Inches(2.4))
for line, sz, col in [("Öğrenilen dersler:", 18, WHITE),
                      ("• Zero-shot yetmez; küçük modelde format/etiket öğrenmek için FT şart.", 15, RGBColor(0xC8, 0xD2, 0xDC)),
                      ("• Dengesizlikte veri-seviyesi tek başına yetersiz; loss-seviyesi gerekli.", 15, RGBColor(0xC8, 0xD2, 0xDC)),
                      ("• Daha büyük model ≠ daha iyi; veri + veriye uyarlanmış mimari belirleyici.", 15, RGBColor(0xC8, 0xD2, 0xDC)),
                      ("• Şeffaf negatif bulgular (cascade, Yeşil F1=0) da bilimsel katkıdır.", 15, RGBColor(0xC8, 0xD2, 0xDC))]:
    p = tf2.add_paragraph(); r = p.add_run(); r.text = line; _set(r, sz, col, bold=(sz == 18))
tf3 = textbox(s, Inches(1.2), SH - Inches(0.9), SW - Inches(2.4), Inches(0.6))
r = tf3.paragraphs[0].add_run()
r.text = "Teşekkürler — Ebubekir Bayar & Furkan Kocataş · github.com/<repo>"
_set(r, 13, RGBColor(0x9F, 0xB3, 0xC8))

out = ROOT / "sunum.pptx"
prs.save(str(out))
print(f"OK -> {out}  ({len(prs.slides)} slayt)")
